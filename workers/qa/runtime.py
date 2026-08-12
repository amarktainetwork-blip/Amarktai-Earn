from __future__ import annotations

import re
import json
import subprocess
from collections import Counter
from html.parser import HTMLParser
from dataclasses import dataclass, field
from pathlib import Path
from typing import Any

from workers.qa.deterministic import verify_csv
from PIL import Image


SUPPORTED_QA_PROFILES = frozenset({
    "csv", "document", "research", "translation", "transcript", "code_patch", "ci", "media",
    "tabular", "spreadsheet", "analysis", "professional_text", "seo_audit", "presentation",
    "produced_document", "public_web", "static_html", "defensive_review", "synthetic_dataset",
    "ai_safety_research",
})


@dataclass(frozen=True)
class QAOutcome:
    passed: bool
    check_type: str
    score: float
    checks: list[str] = field(default_factory=list)
    evidence: dict[str, Any] = field(default_factory=dict)


def _text_file(primary: Path) -> str:
    try:
        return primary.read_text(encoding="utf-8", errors="replace").strip()
    except OSError:
        return ""


def _document_outcome(primary: Path, evidence: dict[str, Any]) -> QAOutcome:
    text = _text_file(primary)
    checks: list[str] = []
    source_chars = int(evidence.get("source_chars") or 0)
    operation = str(evidence.get("operation") or "")
    if text:
        checks.append("output_nonempty")
    if source_chars > 0:
        checks.append("source_text_present")

    passed = bool(text) and source_chars > 0
    if operation == "document_extract_text":
        if source_chars and len(text) >= max(1, int(source_chars * 0.8)):
            checks.append("extraction_length_consistent")
        passed = passed and "extraction_length_consistent" in checks
    elif operation in {"document_summarize", "document_rewrite"}:
        minimum = min(40, max(10, source_chars // 10))
        if len(text) >= minimum:
            checks.append("generated_document_has_content")
        passed = passed and "generated_document_has_content" in checks
    else:
        passed = False

    return QAOutcome(
        passed=passed,
        check_type="document_structural",
        score=1.0 if passed else 0.0,
        checks=checks,
        evidence={"operation": operation, "output_chars": len(text), "source_chars": source_chars},
    )


def _research_outcome(primary: Path, evidence: dict[str, Any]) -> QAOutcome:
    text = _text_file(primary)
    checks: list[str] = []
    sources = [str(url).rstrip(".,;") for url in evidence.get("sources", []) if str(url).startswith("https://")]
    inline = [url.rstrip(".,;") for url in re.findall(r"https://[^\s)\]>]+", text)]
    unique_sources = list(dict.fromkeys([*sources, *inline]))
    if len(text) >= 100:
        checks.append("report_has_substance")
    if len(unique_sources) >= 2:
        checks.append("multiple_https_sources")
    if "sources" in text.casefold() or inline:
        checks.append("citations_present")
    passed = all(name in checks for name in ("report_has_substance", "multiple_https_sources", "citations_present"))
    return QAOutcome(
        passed=passed,
        check_type="research_citation_structural",
        score=1.0 if passed else 0.0,
        checks=checks,
        evidence={"source_count": len(unique_sources), "sources": unique_sources[:20], "output_chars": len(text)},
    )


def _translation_outcome(primary: Path, evidence: dict[str, Any]) -> QAOutcome:
    text = _text_file(primary)
    checks: list[str] = []
    target_language = str(evidence.get("target_language") or "").strip()
    source_chars = int(evidence.get("source_chars") or 0)
    if text:
        checks.append("output_nonempty")
    if target_language:
        checks.append("target_language_explicit")
    if source_chars > 0:
        checks.append("source_text_present")
    if source_chars and text:
        ratio = len(text) / source_chars
        if 0.20 <= ratio <= 5.0:
            checks.append("translation_length_plausible")
    passed = all(name in checks for name in ("output_nonempty", "target_language_explicit", "source_text_present", "translation_length_plausible"))
    return QAOutcome(
        passed=passed,
        check_type="translation_structural",
        score=1.0 if passed else 0.0,
        checks=checks,
        evidence={"target_language": target_language, "source_chars": source_chars, "output_chars": len(text)},
    )


def _transcript_outcome(primary: Path, evidence: dict[str, Any]) -> QAOutcome:
    text = _text_file(primary)
    checks: list[str] = []
    words = len(text.split())
    if text:
        checks.append("output_nonempty")
    if words >= 3:
        checks.append("transcript_has_words")
    passed = "output_nonempty" in checks and "transcript_has_words" in checks
    return QAOutcome(
        passed=passed,
        check_type="transcript_structural",
        score=1.0 if passed else 0.0,
        checks=checks,
        evidence={"word_count": words, "output_chars": len(text)},
    )


def _code_patch_outcome(primary: Path, evidence: dict[str, Any]) -> QAOutcome:
    text = _text_file(primary)
    checks: list[str] = []
    if text and ("diff --git" in text or ("--- " in text and "+++ " in text)):
        checks.append("patch_nonempty")
    if int(evidence.get("agent_exit_code", 1)) == 0:
        checks.append("agent_completed")
    if int(evidence.get("test_exit_code", 1)) == 0:
        checks.append("independent_tests_passed")
    passed = all(name in checks for name in ("patch_nonempty", "agent_completed", "independent_tests_passed"))
    return QAOutcome(
        passed=passed,
        check_type="sandbox_code_patch",
        score=1.0 if passed else 0.0,
        checks=checks,
        evidence={"patch_chars": len(text), "agent_exit_code": evidence.get("agent_exit_code"), "test_exit_code": evidence.get("test_exit_code"), "sandbox_id": evidence.get("sandbox_id")},
    )


def _ci_outcome(primary: Path, evidence: dict[str, Any]) -> QAOutcome:
    text = _text_file(primary)
    checks: list[str] = []
    if text:
        checks.append("test_report_present")
    if int(evidence.get("test_exit_code", 1)) == 0:
        checks.append("tests_passed")
    passed = "test_report_present" in checks and "tests_passed" in checks
    return QAOutcome(
        passed=passed,
        check_type="sandbox_ci",
        score=1.0 if passed else 0.0,
        checks=checks,
        evidence={"report_chars": len(text), "test_exit_code": evidence.get("test_exit_code"), "sandbox_id": evidence.get("sandbox_id")},
    )


def _media_outcome(primary: Path, evidence: dict[str, Any]) -> QAOutcome:
    checks: list[str] = []
    maximum = int(evidence.get("max_output_bytes") or 0)
    if primary.is_file() and primary.stat().st_size > 0:
        checks.append("output_nonempty")
    if maximum and primary.is_file() and primary.stat().st_size <= maximum:
        checks.append("output_size_bounded")
    kind = str(evidence.get("media_kind") or "")
    qa_evidence: dict[str, Any] = {"size_bytes": primary.stat().st_size if primary.is_file() else 0}
    if kind == "image" and primary.is_file():
        try:
            with Image.open(primary) as image:
                image.verify()
            with Image.open(primary) as image:
                dimensions = list(image.size)
                actual_format = str(image.format or "").upper()
            qa_evidence.update({"dimensions": dimensions, "format": actual_format})
            checks.append("image_decodes")
            if actual_format == str(evidence.get("expected_format") or "").upper():
                checks.append("image_format_matches")
            if dimensions == list(evidence.get("expected_dimensions") or []):
                checks.append("image_dimensions_match")
        except (OSError, ValueError, Image.DecompressionBombError, Image.DecompressionBombWarning):
            pass
        required = {"output_nonempty", "output_size_bounded", "image_decodes", "image_format_matches", "image_dimensions_match"}
    elif kind == "av" and primary.is_file():
        try:
            result = subprocess.run(
                ["ffprobe", "-v", "error", "-show_entries", "format=duration,format_name:stream=codec_type", "-of", "json", str(primary)],
                capture_output=True, text=True, timeout=30, check=False,
            )
            data = json.loads(result.stdout) if result.returncode == 0 else {}
            duration = float(data.get("format", {}).get("duration") or 0)
            stream_types = [str(row.get("codec_type")) for row in data.get("streams", [])]
            qa_evidence.update({"duration_seconds": duration, "stream_types": stream_types, "format_name": data.get("format", {}).get("format_name")})
            if duration > 0:
                checks.append("media_probe_valid")
            if evidence.get("require_audio") and "audio" in stream_types:
                checks.append("audio_stream_present")
            if evidence.get("require_video") and "video" in stream_types:
                checks.append("video_stream_present")
            expected_duration = evidence.get("expected_duration_seconds")
            if expected_duration is None or abs(duration - float(expected_duration)) <= 1.0:
                checks.append("duration_within_tolerance")
            expected_format = str(evidence.get("expected_format") or "")
            actual_formats = set(str(data.get("format", {}).get("format_name") or "").split(","))
            matches = (
                (expected_format == "mp4" and "mp4" in actual_formats)
                or (expected_format == "webm" and "webm" in actual_formats)
                or (expected_format == "mp3" and "mp3" in actual_formats)
                or (expected_format == "wav" and "wav" in actual_formats)
            )
            if matches:
                checks.append("media_format_matches")
        except (OSError, subprocess.TimeoutExpired, json.JSONDecodeError, TypeError, ValueError):
            pass
        required = {"output_nonempty", "output_size_bounded", "media_probe_valid", "duration_within_tolerance", "media_format_matches"}
        if evidence.get("require_audio"):
            required.add("audio_stream_present")
        if evidence.get("require_video"):
            required.add("video_stream_present")
    else:
        required = {"unsupported_media_evidence"}
    passed = required.issubset(checks)
    return QAOutcome(passed=passed, check_type="deterministic_media", score=1.0 if passed else 0.0, checks=checks, evidence=qa_evidence)


def _tabular_outcome(primary: Path, evidence: dict[str, Any]) -> QAOutcome:
    checks: list[str] = []
    suffix = primary.suffix.casefold()
    qa_evidence: dict[str, Any] = {"format": suffix.lstrip(".")}
    try:
        if suffix == ".csv":
            result = verify_csv(primary, expected_rows=evidence.get("rows"), required_columns=evidence.get("columns"))
            checks.extend(name for name, passed in result.checks.items() if passed)
            passed = result.passed
            qa_evidence.update(result.evidence)
        elif suffix == ".json":
            payload = json.loads(primary.read_text(encoding="utf-8"))
            if isinstance(payload, list) or (isinstance(payload, dict) and "valid" in payload and "errors" in payload):
                checks.append("json_structure_valid")
            passed = "json_structure_valid" in checks
            qa_evidence["json_type"] = type(payload).__name__
        elif suffix == ".xlsx":
            passed, spreadsheet_checks, spreadsheet_evidence = _inspect_workbook(primary, evidence)
            checks.extend(spreadsheet_checks); qa_evidence.update(spreadsheet_evidence)
        else:
            passed = False
    except Exception:
        passed = False
    return QAOutcome(passed=passed, check_type="deterministic_tabular", score=1.0 if passed else 0.0, checks=checks, evidence=qa_evidence)


def _inspect_workbook(primary: Path, evidence: dict[str, Any]) -> tuple[bool, list[str], dict[str, Any]]:
    from openpyxl import load_workbook

    checks: list[str] = []
    workbook = load_workbook(primary, read_only=True, data_only=False, keep_links=False)
    try:
        sheets = workbook.sheetnames
        if len(sheets) >= int(evidence.get("minimum_sheets") or 1): checks.append("required_sheets_present")
        populated = 0; formulas = 0
        for sheet in workbook.worksheets:
            for row in sheet.iter_rows():
                for cell in row:
                    if cell.value not in (None, ""): populated += 1
                    if cell.data_type == "f": formulas += 1
        if populated: checks.append("workbook_has_content")
        expected_formulas = int(evidence.get("generated_formula_cells") or 0)
        if formulas == expected_formulas and evidence.get("formula_injection_neutralized") is True: checks.append("formula_cells_authorized")
        expected_names = set(evidence.get("sheet_names") or [])
        if not expected_names or expected_names.issubset(sheets): checks.append("sheet_names_match")
        required = {"required_sheets_present", "workbook_has_content", "formula_cells_authorized", "sheet_names_match"}
        return required.issubset(checks), checks, {"sheet_names": sheets, "populated_cells": populated, "formula_cells": formulas}
    finally:
        workbook.close()


def _spreadsheet_outcome(primary: Path, evidence: dict[str, Any], *, analysis: bool = False) -> QAOutcome:
    try:
        passed, checks, qa_evidence = _inspect_workbook(primary, evidence)
    except Exception:
        passed, checks, qa_evidence = False, [], {}
    if analysis:
        if evidence.get("regulated_advice") is False: checks.append("not_regulated_advice")
        if isinstance(evidence.get("visualization_present"), bool): checks.append("visualization_state_declared")
        passed = passed and {"not_regulated_advice", "visualization_state_declared"}.issubset(checks)
    return QAOutcome(passed=passed, check_type="deterministic_analysis" if analysis else "deterministic_spreadsheet", score=1.0 if passed else 0.0, checks=checks, evidence=qa_evidence)


def _professional_text_outcome(primary: Path, evidence: dict[str, Any]) -> QAOutcome:
    text = _text_file(primary); checks = []
    if len(text) >= 80: checks.append("deliverable_has_substance")
    if evidence.get("deliverable_type"): checks.append("deliverable_type_explicit")
    if evidence.get("operation") in {"technical_documentation", "content_package", "support_content_package"}: checks.append("operation_supported")
    if evidence.get("operation") != "support_content_package" or evidence.get("draft_only") is True: checks.append("external_send_not_claimed")
    passed = {"deliverable_has_substance", "deliverable_type_explicit", "operation_supported", "external_send_not_claimed"}.issubset(checks)
    return QAOutcome(passed=passed, check_type="professional_text_structural", score=1.0 if passed else 0.0, checks=checks, evidence={"output_chars": len(text), "deliverable_type": evidence.get("deliverable_type")})


def _seo_outcome(primary: Path, evidence: dict[str, Any]) -> QAOutcome:
    checks = []; payload = {}
    try: payload = json.loads(primary.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError): pass
    required_keys = {"word_count", "headings", "top_keywords", "findings", "policy"}
    if required_keys.issubset(payload): checks.append("structured_audit_present")
    if payload.get("policy", {}).get("spam_or_backlink_automation") is False: checks.append("spam_automation_absent")
    passed = {"structured_audit_present", "spam_automation_absent"}.issubset(checks)
    return QAOutcome(passed=passed, check_type="deterministic_seo_audit", score=1.0 if passed else 0.0, checks=checks, evidence={"finding_count": len(payload.get("findings", [])) if isinstance(payload, dict) else 0})


def _presentation_outcome(primary: Path, evidence: dict[str, Any]) -> QAOutcome:
    from pptx import Presentation

    checks = []; slide_count = 0; titled = 0
    try:
        deck = Presentation(primary); slide_count = len(deck.slides)
        titled = sum(bool(slide.shapes.title and slide.shapes.title.text.strip()) for slide in deck.slides)
        if slide_count == int(evidence.get("required_slide_count") or 0) and slide_count > 0: checks.append("slide_count_matches")
        if titled == slide_count: checks.append("all_slides_titled")
        if evidence.get("reopened") is True: checks.append("worker_reopened_deck")
    except Exception: pass
    passed = {"slide_count_matches", "all_slides_titled", "worker_reopened_deck"}.issubset(checks)
    return QAOutcome(passed=passed, check_type="deterministic_presentation", score=1.0 if passed else 0.0, checks=checks, evidence={"slide_count": slide_count, "titled_slides": titled})


def _produced_document_outcome(primary: Path, evidence: dict[str, Any]) -> QAOutcome:
    checks = []; qa_evidence = {"format": primary.suffix.casefold().lstrip(".")}
    try:
        if primary.suffix.casefold() == ".docx":
            from docx import Document
            document = Document(primary); count = sum(bool(row.text.strip()) for row in document.paragraphs)
            qa_evidence["paragraph_count"] = count
            if count >= 2: checks.append("docx_has_structure")
        elif primary.suffix.casefold() == ".pdf":
            from pypdf import PdfReader
            reader = PdfReader(str(primary)); count = len(reader.pages)
            qa_evidence["page_count"] = count
            if count >= 1 and any((page.extract_text() or "").strip() for page in reader.pages): checks.append("pdf_has_extractable_content")
        if evidence.get("reopened") is True: checks.append("worker_reopened_document")
    except Exception: pass
    passed = "worker_reopened_document" in checks and any(name in checks for name in ("docx_has_structure", "pdf_has_extractable_content"))
    return QAOutcome(passed=passed, check_type="deterministic_produced_document", score=1.0 if passed else 0.0, checks=checks, evidence=qa_evidence)


def _public_web_outcome(primary: Path, evidence: dict[str, Any]) -> QAOutcome:
    checks = []; payload = {}
    try: payload = json.loads(primary.read_text(encoding="utf-8"))
    except (OSError, json.JSONDecodeError): pass
    if str(payload.get("url") or "").startswith("https://"): checks.append("https_final_url")
    if payload.get("text") and payload.get("purpose"): checks.append("bounded_extract_present")
    if evidence.get("robots_checked") is True: checks.append("robots_checked")
    if evidence.get("policy_confirmed") is True: checks.append("policy_confirmed")
    passed = {"https_final_url", "bounded_extract_present", "robots_checked", "policy_confirmed"}.issubset(checks)
    return QAOutcome(passed=passed, check_type="bounded_public_web", score=1.0 if passed else 0.0, checks=checks, evidence={"bytes": payload.get("bytes"), "url": payload.get("url")})


class _HTMLContractParser(HTMLParser):
    def __init__(self):
        super().__init__(); self.lang = ""; self.title = False; self.h1 = 0; self.scripts = 0; self.forms = 0; self.skip = False
    def handle_starttag(self, tag, attrs):
        values = dict(attrs); tag = tag.casefold()
        if tag == "html": self.lang = values.get("lang", "")
        if tag == "title": self.title = True
        if tag == "h1": self.h1 += 1
        if tag == "script": self.scripts += 1
        if tag == "form": self.forms += 1
        if tag == "a" and values.get("href") == "#main": self.skip = True


def _static_html_outcome(primary: Path, evidence: dict[str, Any]) -> QAOutcome:
    checks = []; parser = _HTMLContractParser()
    try: parser.feed(primary.read_text(encoding="utf-8"))
    except OSError: pass
    if parser.lang: checks.append("language_declared")
    if parser.title and parser.h1 == 1: checks.append("title_and_single_h1")
    if parser.skip: checks.append("skip_link_present")
    if parser.scripts == 0 and parser.forms == 0: checks.append("no_active_or_sending_surface")
    if evidence.get("deployment_performed") is False: checks.append("deployment_not_claimed")
    passed = {"language_declared", "title_and_single_h1", "skip_link_present", "no_active_or_sending_surface", "deployment_not_claimed"}.issubset(checks)
    return QAOutcome(passed=passed, check_type="deterministic_static_html", score=1.0 if passed else 0.0, checks=checks, evidence={"h1_count": parser.h1, "scripts": parser.scripts, "forms": parser.forms})


def _defensive_review_outcome(primary: Path, evidence: dict[str, Any]) -> QAOutcome:
    text = _text_file(primary); checks = []
    if len(text) >= 150 and "Defensive Code Review" in text: checks.append("review_report_present")
    if evidence.get("authorization_confirmed") is True and evidence.get("scope"): checks.append("authorization_and_scope_recorded")
    if evidence.get("network_testing_performed") is False and evidence.get("exploitation_performed") is False: checks.append("defensive_boundary_preserved")
    passed = {"review_report_present", "authorization_and_scope_recorded", "defensive_boundary_preserved"}.issubset(checks)
    return QAOutcome(passed=passed, check_type="deterministic_defensive_review", score=1.0 if passed else 0.0, checks=checks, evidence={"files_scanned": evidence.get("files_scanned"), "finding_count": evidence.get("finding_count")})


def _synthetic_dataset_outcome(primary: Path, evidence: dict[str, Any]) -> QAOutcome:
    checks = []
    rows = []
    try:
        rows = [json.loads(line) for line in primary.read_text(encoding="utf-8").splitlines() if line.strip()]
    except (OSError, json.JSONDecodeError):
        rows = []
    schema = evidence.get("schema") if isinstance(evidence.get("schema"), dict) else {}
    fields = schema.get("fields") if isinstance(schema.get("fields"), dict) else {}
    if rows and len(rows) == int(evidence.get("accepted_records") or -1):
        checks.append("accepted_count_reopened")
    if rows and all(set(row) == {*fields, "_split"} for row in rows if isinstance(row, dict)):
        checks.append("schema_fields_exact")
    canonical = [json.dumps({key: value for key, value in row.items() if key != "_split"}, sort_keys=True, ensure_ascii=False) for row in rows]
    if canonical and len(canonical) == len(set(canonical)):
        checks.append("no_cross_split_duplicates")
    split_counts = Counter(str(row.get("_split")) for row in rows if isinstance(row, dict))
    if dict(split_counts) == evidence.get("split_counts"):
        checks.append("split_distribution_reopened")
    card_path = Path(str(evidence.get("dataset_card_path") or ""))
    try:
        card_text = card_path.read_text(encoding="utf-8")
    except OSError:
        card_text = ""
    if "# Synthetic Dataset Card" in card_text and "Privacy/provenance" in card_text:
        checks.append("dataset_card_present")
    if evidence.get("rights_confirmed") is True and evidence.get("provenance"):
        checks.append("rights_and_provenance_recorded")
    if int(evidence.get("pii_rejected") or 0) >= 0 and int(evidence.get("contamination_rejected") or 0) >= 0:
        checks.append("privacy_and_contamination_metrics_recorded")
    required = {
        "accepted_count_reopened", "schema_fields_exact", "no_cross_split_duplicates",
        "split_distribution_reopened", "dataset_card_present", "rights_and_provenance_recorded",
        "privacy_and_contamination_metrics_recorded",
    }
    passed = required.issubset(checks)
    return QAOutcome(
        passed=passed, check_type="deterministic_synthetic_dataset", score=1.0 if passed else 0.0,
        checks=checks, evidence={"reopened_records": len(rows), "split_counts": dict(split_counts)},
    )


def _ai_safety_outcome(primary: Path, evidence: dict[str, Any]) -> QAOutcome:
    text = _text_file(primary)
    checks = []
    if "# Authorized AI Safety Research Report" in text and len(text) >= 300:
        checks.append("professional_report_present")
    if evidence.get("authorization_hash") and evidence.get("scope_version_id") and evidence.get("target_id"):
        checks.append("persisted_authorization_boundary")
    executed = int(evidence.get("requests_executed") or 0)
    maximum = int(evidence.get("max_requests") or 0)
    if 0 < executed <= maximum and int(evidence.get("rate_limit_per_minute") or 0) > 0:
        checks.append("request_and_rate_bounds_recorded")
    if evidence.get("network_testing_performed") is False and evidence.get("remote_target_interaction") is False:
        checks.append("offline_fixture_boundary_preserved")
    if evidence.get("raw_prompts_in_artifact") is False and evidence.get("private_target_data_in_artifact") is False:
        checks.append("sensitive_and_harmful_details_excluded")
    if "independent reproduction" in text.casefold() and "duplicate check" in text.casefold():
        checks.append("submission_gate_documented")
    required = {
        "professional_report_present", "persisted_authorization_boundary", "request_and_rate_bounds_recorded",
        "offline_fixture_boundary_preserved", "sensitive_and_harmful_details_excluded", "submission_gate_documented",
    }
    passed = required.issubset(checks)
    return QAOutcome(
        passed=passed, check_type="deterministic_ai_safety_research", score=1.0 if passed else 0.0,
        checks=checks, evidence={"requests_executed": executed, "candidate_findings": len(evidence.get("findings") or [])},
    )


def run_qa(profile: str, primary: Path, worker_evidence: dict[str, Any] | None = None) -> QAOutcome:
    evidence = worker_evidence if isinstance(worker_evidence, dict) else {}
    if profile == "csv":
        result = verify_csv(primary, expected_rows=evidence.get("rows"), required_columns=evidence.get("columns"))
        return QAOutcome(
            passed=result.passed,
            check_type="deterministic_csv",
            score=1.0 if result.passed else 0.0,
            checks=list(result.checks),
            evidence=dict(result.evidence),
        )
    if profile == "document":
        return _document_outcome(primary, evidence)
    if profile == "research":
        return _research_outcome(primary, evidence)
    if profile == "translation":
        return _translation_outcome(primary, evidence)
    if profile == "transcript":
        return _transcript_outcome(primary, evidence)
    if profile == "code_patch":
        return _code_patch_outcome(primary, evidence)
    if profile == "ci":
        return _ci_outcome(primary, evidence)
    if profile == "media":
        return _media_outcome(primary, evidence)
    if profile == "tabular":
        return _tabular_outcome(primary, evidence)
    if profile == "spreadsheet":
        return _spreadsheet_outcome(primary, evidence)
    if profile == "analysis":
        return _spreadsheet_outcome(primary, evidence, analysis=True)
    if profile == "professional_text":
        return _professional_text_outcome(primary, evidence)
    if profile == "seo_audit":
        return _seo_outcome(primary, evidence)
    if profile == "presentation":
        return _presentation_outcome(primary, evidence)
    if profile == "produced_document":
        return _produced_document_outcome(primary, evidence)
    if profile == "public_web":
        return _public_web_outcome(primary, evidence)
    if profile == "static_html":
        return _static_html_outcome(primary, evidence)
    if profile == "defensive_review":
        return _defensive_review_outcome(primary, evidence)
    if profile == "synthetic_dataset":
        return _synthetic_dataset_outcome(primary, evidence)
    if profile == "ai_safety_research":
        return _ai_safety_outcome(primary, evidence)
    raise ValueError(f"unsupported QA profile: {profile}")
