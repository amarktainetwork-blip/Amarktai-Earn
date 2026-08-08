from __future__ import annotations

import os
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt

from workers.base import WorkRequest, WorkResult, Worker


def _slides(inputs: dict) -> list[dict[str, str]]:
    supplied = inputs.get("slides")
    if isinstance(supplied, list):
        slides = [{"title": str(row.get("title") or "Untitled"), "body": str(row.get("body") or "")} for row in supplied if isinstance(row, dict)]
    else:
        source = inputs.get("source")
        text = Path(str(source)).read_text(encoding="utf-8", errors="replace") if source else str(inputs.get("content") or inputs.get("brief") or "")
        slides = []
        current = None
        for line in text.splitlines():
            if line.lstrip().startswith("#"):
                if current: slides.append(current)
                current = {"title": line.lstrip("# ").strip() or "Untitled", "body": ""}
            elif current and line.strip():
                current["body"] += ("\n" if current["body"] else "") + line.strip().lstrip("-* ")
        if current: slides.append(current)
    if not slides:
        slides = [{"title": str(inputs.get("title") or "Presentation"), "body": str(inputs.get("brief") or "Content supplied by the customer.")}]
    maximum = max(1, min(int(os.getenv("PRESENTATION_MAX_SLIDES", "30")), 100))
    if len(slides) > maximum:
        raise ValueError("PRESENTATION_SLIDE_LIMIT")
    return slides


class PresentationWorker(Worker):
    worker_class = "presentations"

    def execute(self, request: WorkRequest) -> WorkResult:
        try:
            if request.inputs.get("operation") != "presentation_create":
                return WorkResult(ok=False, error="unsupported presentation operation")
            rows = _slides(request.inputs)
            request.workspace.mkdir(parents=True, exist_ok=True)
            target = request.workspace / "presentation.pptx"
            deck = Presentation()
            deck.slide_width = Inches(13.333); deck.slide_height = Inches(7.5)
            for index, row in enumerate(rows):
                layout = deck.slide_layouts[0] if index == 0 else deck.slide_layouts[1]
                slide = deck.slides.add_slide(layout)
                slide.shapes.title.text = row["title"][:200]
                if len(slide.placeholders) > 1:
                    frame = slide.placeholders[1].text_frame
                    frame.clear()
                    lines = [line.strip() for line in row["body"].splitlines() if line.strip()] or [" "]
                    for line_index, line in enumerate(lines[:12]):
                        paragraph = frame.paragraphs[0] if line_index == 0 else frame.add_paragraph()
                        paragraph.text = line[:500]; paragraph.font.size = Pt(22)
            deck.core_properties.title = str(request.inputs.get("title") or rows[0]["title"])
            deck.save(target)
            reopened = Presentation(target)
            titles = [slide.shapes.title.text.strip() for slide in reopened.slides if slide.shapes.title]
            return WorkResult(ok=True, artifacts=[target], evidence={
                "operation": "presentation_create", "slide_count": len(reopened.slides),
                "titles": titles, "required_slide_count": len(rows), "reopened": True,
            })
        except (OSError, KeyError, TypeError, ValueError) as exc:
            return WorkResult(ok=False, error=str(exc))
